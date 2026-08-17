from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("Smart_E_Home_Presentation.pptx")
BG = RGBColor(8, 18, 32)
CARD = RGBColor(18, 35, 55)
CYAN = RGBColor(61, 214, 208)
GREEN = RGBColor(104, 225, 160)
WHITE = RGBColor(245, 249, 255)
MUTED = RGBColor(167, 184, 204)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def text(slide, value, x, y, w, h, size=24, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def base(title, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(.12), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    text(slide, title, .7, .45, 11.8, .65, 28, WHITE, True)
    text(slide, f"SMART E-HOME  •  {number:02d}", .72, 7.05, 3.2, .25, 9, MUTED)
    return slide


def card(slide, title, body, x, y, w, h, accent=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = accent
    text(slide, title, x + .25, y + .2, w - .5, .4, 16, accent, True)
    text(slide, body, x + .25, y + .72, w - .5, h - .85, 15, WHITE)


# 1 — title
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(.14), prs.slide_height)
accent.fill.solid()
accent.fill.fore_color.rgb = CYAN
text(s, "SMART E-HOME", .8, 1.35, 8.0, .9, 42, WHITE, True)
text(s, "ESP32-based IoT home monitoring and control", .82, 2.35, 8.5, .6, 23, CYAN)
text(s, "Presented by: [Your Name]", .82, 3.25, 5.5, .45, 18, MUTED)
text(s, "Final Project Presentation", .82, 3.78, 5.5, .4, 16, MUTED)
hero = Path("src/assets/hero.png")
if hero.exists():
    s.shapes.add_picture(str(hero), Inches(8.8), Inches(1.0), width=Inches(3.8))
text(s, "Opening: What if we could monitor and control home devices from one live dashboard?", .82, 6.65, 11.7, .4, 13, GREEN)

# 2
s = base("The Problem", 2)
card(s, "Scattered controls", "Home devices are controlled separately, making everyday management less convenient.", .8, 1.45, 3.7, 2.1)
card(s, "Limited visibility", "Users may not know the current temperature, humidity, gas status, or device state.", 4.8, 1.45, 3.7, 2.1)
card(s, "Slow response", "Without a connected dashboard, reacting to an unsafe or unusual condition can take longer.", 8.8, 1.45, 3.7, 2.1)
text(s, "Goal: build one simple interface for live monitoring and reliable device control.", 1.3, 4.55, 10.7, .7, 24, GREEN, True, PP_ALIGN.CENTER)

# 3
s = base("Our Solution", 3)
text(s, "A web dashboard connected to an ESP32 through MQTT.", .85, 1.25, 11.6, .55, 25, CYAN, True)
card(s, "Monitor", "View live temperature, humidity, gas reading, MQTT status, and ESP32 connectivity.", .9, 2.25, 3.6, 2.25)
card(s, "Control", "Turn five connected devices on or off individually, or use all-on and all-off controls.", 4.85, 2.25, 3.6, 2.25)
card(s, "Record", "Save sensor readings and device-function logs in PostgreSQL for later review.", 8.8, 2.25, 3.6, 2.25)

# 4
s = base("How the System Works", 4)
items = [("React UI", .65), ("Express API", 3.2), ("MQTT Broker", 5.75), ("ESP32", 8.3), ("Devices + Sensors", 10.85)]
for label, x in items:
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.45), Inches(1.85), Inches(1.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD; shape.line.color.rgb = CYAN
    text(s, label, x + .1, 2.77, 1.65, .35, 15, WHITE, True, PP_ALIGN.CENTER)
for x in [2.55, 5.1, 7.65, 10.2]:
    text(s, "→", x, 2.7, .5, .4, 25, GREEN, True, PP_ALIGN.CENTER)
text(s, "Commands flow to the hardware; device states and sensor data flow back to the dashboard.", 1.0, 4.35, 11.2, .75, 19, MUTED, False, PP_ALIGN.CENTER)
text(s, "PostgreSQL stores sensor logs and function logs.", 1.0, 5.25, 11.2, .5, 18, GREEN, True, PP_ALIGN.CENTER)

# 5
s = base("Main Features", 5)
features = [
    ("Demo login", "Validates email and password before opening the dashboard."),
    ("Device control", "Controls 3 lights, a buzzer, and a fan."),
    ("Group actions", "Turn all devices on or off from one place."),
    ("Live sensors", "Displays temperature, humidity, and gas state."),
    ("Status checks", "Shows API, MQTT, and ESP32 connectivity."),
    ("Error handling", "Explains offline backend, broker, or ESP32 conditions."),
]
for i, (title, body) in enumerate(features):
    x = .75 + (i % 3) * 4.15
    y = 1.35 + (i // 3) * 2.45
    card(s, title, body, x, y, 3.75, 1.9, GREEN if i % 2 else CYAN)

# 6
s = base("Technology Stack", 6)
stack = [
    ("Frontend", "React 19 • Vite • CSS • Lucide icons"),
    ("Backend", "Node.js • Express • REST API"),
    ("IoT communication", "MQTT • Aedes local broker • mqtt.js"),
    ("Hardware", "ESP32 • connected devices • environmental sensors"),
    ("Database", "PostgreSQL • sensor logs • function logs"),
    ("Deployment", "Frontend API URL configured with environment variables"),
]
for i, (title, body) in enumerate(stack):
    x = .85 + (i % 2) * 6.2
    y = 1.25 + (i // 2) * 1.65
    card(s, title, body, x, y, 5.75, 1.35)

# 7
s = base("Demo Flow", 7)
steps = ["1  Login", "2  Check API / MQTT / ESP32", "3  View sensor values", "4  Toggle one device", "5  Use Turn All On / Off", "6  Explain offline safety message"]
for i, item in enumerate(steps):
    y = 1.2 + i * .83
    color = CYAN if i < 3 else GREEN
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(.9), Inches(y), Inches(.28), Inches(.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    text(s, item, 1.45, y - .05, 10.8, .42, 20, WHITE, i == 0)
text(s, "Tip: keep the live demo under 2 minutes.", 8.25, 6.45, 4.2, .4, 14, MUTED, False, PP_ALIGN.RIGHT)

# 8
s = base("Challenges and Learning", 8)
card(s, "Keeping states synchronized", "The browser, server, MQTT broker, and ESP32 must agree on each device state.", .9, 1.35, 3.7, 2.2)
card(s, "Handling disconnections", "The app checks whether MQTT and ESP32 are online before sending commands.", 4.82, 1.35, 3.7, 2.2)
card(s, "Connecting software to hardware", "MQTT topics provide a clear message path between the web API and ESP32.", 8.75, 1.35, 3.7, 2.2)
text(s, "My biggest learning: [Replace this with your own honest learning in one sentence.]", 1.0, 4.65, 11.3, .8, 21, GREEN, True, PP_ALIGN.CENTER)

# 9
s = base("Limitations and Future Work", 9)
card(s, "Current limitations", "• Demo-only login\n• Device list is stored in server memory\n• Dashboard does not yet display saved logs\n• Hardware availability affects the live demo", .85, 1.35, 5.75, 3.8, CYAN)
card(s, "Possible future work", "• Secure authentication\n• User-specific homes and devices\n• Log-history charts\n• Alerts and automation rules\n• Improved production MQTT hosting", 6.75, 1.35, 5.75, 3.8, GREEN)
text(s, "Future work is separate from the current MVP.", .95, 5.75, 11.4, .5, 17, MUTED, False, PP_ALIGN.CENTER)

# 10
s = base("Thank You", 10)
text(s, "Questions?", 1.0, 2.0, 11.3, 1.0, 42, WHITE, True, PP_ALIGN.CENTER)
text(s, "SMART E-HOME", 1.0, 3.25, 11.3, .7, 25, CYAN, True, PP_ALIGN.CENTER)
text(s, "A connected dashboard for safer and simpler home control.", 1.0, 4.05, 11.3, .55, 18, MUTED, False, PP_ALIGN.CENTER)
text(s, "Before presenting: replace [Your Name], add one dashboard screenshot, and personalize the learning slide.", 1.0, 6.2, 11.3, .5, 14, GREEN, False, PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Created {OUT.resolve()} with {len(prs.slides)} slides")
